import os
import shutil
import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def build_presentation():
    input_path = "ppt/SIH-official-template.pptx"
    output_path = "ppt/SIH-official-template.pptx"
    clean_copy_path = "ppt/SurakshaMesh-SIH26025.pptx"
    downloads_path = os.path.expanduser("~/Downloads/SurakshaMesh-SIH26025.pptx")

    prs = pptx.Presentation(input_path)

    # Color Palette
    DARK_BLUE = RGBColor(15, 30, 60)
    ACCENT_BLUE = RGBColor(30, 90, 200)
    TEXT_DARK = RGBColor(20, 20, 20)
    TEXT_MUTED = RGBColor(80, 80, 80)
    ALERT_RED = RGBColor(200, 30, 30)

    # ==========================================
    # SLIDE 1: TITLE PAGE
    # ==========================================
    slide1 = prs.slides[0]
    for shape in slide1.shapes:
        if shape.name == "Title 7":
            shape.text_frame.text = "SMART INDIA HACKATHON 2026"
            shape.text_frame.paragraphs[0].font.size = Pt(28)
            shape.text_frame.paragraphs[0].font.bold = True
        elif shape.name == "Subtitle 3":
            shape.text_frame.clear()
            p = shape.text_frame.paragraphs[0]
            p.text = "SurakshaMesh — Low-Cost Wireless Surface Mesh for Real-Time Mine Subsidence Early Warning"
            p.font.size = Pt(17)
            p.font.bold = True
            p.font.color.rgb = ACCENT_BLUE
        elif shape.name == "TextBox 9":
            shape.text_frame.clear()
            tf = shape.text_frame
            tf.word_wrap = True
            
            lines = [
                ("Problem Statement ID:", " SIH26025"),
                ("Problem Statement Title:", " Development of an AI-enabled Low Cost Real Time Mine Subsidence Monitoring, Prediction and Early Warning System for Underground Coal Mines in India"),
                ("Organization:", " Coal India Limited (CIL) / Ministry of Coal"),
                ("Theme:", " Disaster Management"),
                ("PS Category:", " Hardware"),
                ("Team Name:", " SurakshaMesh (Team DhartiNode)")
            ]
            for label_text, val_text in lines:
                p = tf.add_paragraph()
                p.space_after = Pt(8)
                run1 = p.add_run()
                run1.text = label_text
                run1.font.bold = True
                run1.font.size = Pt(13)
                run1.font.color.rgb = DARK_BLUE

                run2 = p.add_run()
                run2.text = val_text
                run2.font.bold = False
                run2.font.size = Pt(13)
                run2.font.color.rgb = TEXT_DARK

    # Update Team Name in Header Ovals for all slides
    for slide in prs.slides:
        for shape in slide.shapes:
            if "Oval" in shape.name or shape.name in ["Oval 8", "Oval 9", "Oval 10", "Oval 11"]:
                shape.text_frame.text = "SurakshaMesh"
                for p in shape.text_frame.paragraphs:
                    p.font.size = Pt(11)
                    p.font.bold = True

    # ==========================================
    # SLIDE 2: IDEA TITLE & PROPOSED SOLUTION
    # ==========================================
    slide2 = prs.slides[1]
    for shape in slide2.shapes:
        if shape.name == "Title 1":
            shape.text_frame.text = "IDEA TITLE: SurakshaMesh — Surface Mesh, First Alert"
            shape.text_frame.paragraphs[0].font.size = Pt(24)
            shape.text_frame.paragraphs[0].font.bold = True
        elif shape.name == "TextBox 8":
            shape.text_frame.clear()
            tf = shape.text_frame
            tf.word_wrap = True

            def add_sec(title, bullets):
                p_head = tf.add_paragraph()
                p_head.space_before = Pt(4)
                p_head.space_after = Pt(2)
                run = p_head.add_run()
                run.text = title
                run.font.bold = True
                run.font.size = Pt(13)
                run.font.color.rgb = ACCENT_BLUE
                for b in bullets:
                    p = tf.add_paragraph()
                    p.level = 0
                    p.space_after = Pt(2)
                    run_b = p.add_run()
                    run_b.text = "• " + b
                    run_b.font.size = Pt(11.5)
                    run_b.font.color.rgb = TEXT_DARK

            add_sec("1. Proposed Solution & Architecture", [
                "Autonomous wireless sensor nodes arranged in a surface mesh directly above underground coal extraction panels.",
                "Edge Sensor Fusion: 6-Axis MEMS IMU (pitch, roll @ 50 Hz) + rolling Vibration RMS + local Active Buzzer & LED alarms.",
                "Zero-Latency Mesh: ESP-NOW peer-to-peer radio protocol -> Central Gateway -> Real-time Hexagonal Digital Twin Dashboard.",
                "Working Prototype: 3 physical hardware nodes with auto-tare calibration, multi-node spatial coherence, and scheduled blast filter."
            ])
            add_sec("2. How It Solves Coal India's Problem", [
                "Current method: Manual monthly DGPS surveys detect subsidence post-facto (after surface houses & haul roads have cracked).",
                "SurakshaMesh: Always-on 24/7 continuous micro-movement tracking as a first-alert layer; DGPS remains periodic calibration."
            ])
            add_sec("3. Innovation and Uniqueness", [
                "Multi-Node Spatial Coherence: Distinguishes large-scale geological subsidence from single-point vehicle or animal noise.",
                "Ultra-Low Cost: ₹1,500 benchtop node / ₹6,000–8,000 industrial solar field station vs ₹15+ Lakhs imported surveying gear.",
                "Offline-First & Blast-Aware: Runs standalone with local audio sirens even if internet dies; includes a 60s scheduled blast mute window."
            ])

    # ==========================================
    # SLIDE 3: TECHNICAL APPROACH
    # ==========================================
    slide3 = prs.slides[2]
    for shape in slide3.shapes:
        if shape.name == "Title 1":
            shape.text_frame.text = "TECHNICAL APPROACH & SYSTEM ARCHITECTURE"
            shape.text_frame.paragraphs[0].font.size = Pt(24)
            shape.text_frame.paragraphs[0].font.bold = True
        elif shape.name == "TextBox 8":
            shape.text_frame.clear()
            tf = shape.text_frame
            tf.word_wrap = True

            def add_sec_3(title, bullets):
                p_head = tf.add_paragraph()
                p_head.space_before = Pt(4)
                p_head.space_after = Pt(2)
                run = p_head.add_run()
                run.text = title
                run.font.bold = True
                run.font.size = Pt(13)
                run.font.color.rgb = ACCENT_BLUE
                for b in bullets:
                    p = tf.add_paragraph()
                    p.space_after = Pt(2)
                    run_b = p.add_run()
                    run_b.text = "• " + b
                    run_b.font.size = Pt(11.5)
                    run_b.font.color.rgb = TEXT_DARK

            add_sec_3("1. Technologies & Hardware Stack", [
                "Hardware: ESP32 (Xtensa Dual-Core 240MHz), MPU6050 6-DOF IMU, TP4056 + 21700 LiPo (4000mAh, 3.7V) battery, Active Buzzer.",
                "Firmware: Embedded C/C++ on FreeRTOS, native I2C register driver (400kHz), Complementary Filter (98% Gyro + 2% Accel @ 50Hz).",
                "Software: Next.js 16 (App Router), React 19, TypeScript, Tailwind CSS, Node.js Serial-to-HTTP Gateway Bridge, Web Audio API."
            ])
            add_sec_3("2. AI & Geotechnical Decision Engine", [
                "Dynamic Ground-Zero Tare: Calibrates initial resting angles to evaluate pure relative ground deflection (ΔPitch, ΔRoll).",
                "Statistical Baseline Learning: Welford's algorithm accumulates 60 quiet samples to learn the ambient noise floor (mean & std dev).",
                "Unsupervised Anomaly Detection: Isolation Forest on multi-variate feature vector [ΔPitch, ΔRoll, PitchRate, RollRate, VibRMS, STA/LTA].",
                "Spatial Coherence Consensus: Confirms active strata sag when 2+ adjacent nodes tilt in the same direction simultaneously."
            ])
            add_sec_3("3. Data Flow Pipeline", [
                "Surface Nodes (50Hz IMU) ──[ESP-NOW Mesh]──> Central Gateway ──[Serial Bridge]──> Dashboard & AI Engine ──> Live 3D Hex Map & Audio Siren"
            ])

    # ==========================================
    # SLIDE 4: FEASIBILITY AND VIABILITY
    # ==========================================
    slide4 = prs.slides[3]
    for shape in slide4.shapes:
        if shape.name == "Title 1":
            shape.text_frame.text = "FEASIBILITY, VIABILITY & RISK MITIGATION"
            shape.text_frame.paragraphs[0].font.size = Pt(24)
            shape.text_frame.paragraphs[0].font.bold = True
        elif shape.name == "TextBox 8":
            shape.text_frame.clear()
            tf = shape.text_frame
            tf.word_wrap = True

            def add_sec_4(title, bullets):
                p_head = tf.add_paragraph()
                p_head.space_before = Pt(4)
                p_head.space_after = Pt(2)
                run = p_head.add_run()
                run.text = title
                run.font.bold = True
                run.font.size = Pt(13)
                run.font.color.rgb = ACCENT_BLUE
                for b in bullets:
                    p = tf.add_paragraph()
                    p.space_after = Pt(2)
                    run_b = p.add_run()
                    run_b.text = "• " + b
                    run_b.font.size = Pt(11.5)
                    run_b.font.color.rgb = TEXT_DARK

            add_sec_4("1. Feasibility & Commercial Viability", [
                "100% COTS Components: Uses commercially available, ruggedized hardware (ESP32, MEMS, Li-ion, sub-GHz radio).",
                "Live Working 3-Node Prototype: Built and verified with live wireless mesh telemetry, auto-tare calibration, and audio alarms.",
                "High Cost Feasibility: Field station at ~₹6,000–8,000 -> Entire 8-node panel mesh demo under ₹60,000 (cheaper than 1 manual survey day)."
            ])
            add_sec_4("2. Key Challenges & Engineering Solutions", [
                "MEMS Thermal Drift: Temperature swings (15°C–45°C) cause apparent tilt -> Solved via onboard MPU6050 temperature logging.",
                "Mine Blasting False Alarms: Heavy dynamite shocks -> Solved via 3s persistence filter + One-Click 60s Scheduled Blast Mode.",
                "Radio Range in Open Mines: 2.4GHz Wi-Fi blocked by spoil heaps -> Transition to Sub-GHz LoRa (SX1278 433/868 MHz) for 2–5 km range.",
                "Weather & Coal Dust: Ingress risks -> IP65 weather-sealed black box enclosure with rubber gaskets and solar top lid."
            ])

    # ==========================================
    # SLIDE 5: IMPACT AND BENEFITS
    # ==========================================
    slide5 = prs.slides[4]
    for shape in slide5.shapes:
        if shape.name == "Title 1":
            shape.text_frame.text = "IMPACT, BENEFITS & SOCIAL VALUE"
            shape.text_frame.paragraphs[0].font.size = Pt(24)
            shape.text_frame.paragraphs[0].font.bold = True
        elif shape.name == "TextBox 8":
            shape.text_frame.clear()
            tf = shape.text_frame
            tf.word_wrap = True

            def add_sec_5(title, bullets):
                p_head = tf.add_paragraph()
                p_head.space_before = Pt(4)
                p_head.space_after = Pt(2)
                run = p_head.add_run()
                run.text = title
                run.font.bold = True
                run.font.size = Pt(13)
                run.font.color.rgb = ACCENT_BLUE
                for b in bullets:
                    p = tf.add_paragraph()
                    p.space_after = Pt(2)
                    run_b = p.add_run()
                    run_b.text = "• " + b
                    run_b.font.size = Pt(11.5)
                    run_b.font.color.rgb = TEXT_DARK

            add_sec_5("1. Target Audience & Stakeholders", [
                "Coal India Limited (CIL) colliery managers, CMPDI survey teams, Directorate General of Mines Safety (DGMS).",
                "Communities & infrastructure located above Bord-and-Pillar / Longwall extraction panels (Jharia, Raniganj, Korba coalfields)."
            ])
            add_sec_5("2. Key Benefits & Real-World Value", [
                "Life Safety (Social): Early evacuation warning hours-to-days before sudden sinkhole collapse, preventing loss of life.",
                "Operational Continuity (Economic): Prevents emergency haul road blockages, equipment burial, and multi-crore structural damage.",
                "Disaster Containment (Environmental): Early detection of surface fissure openings prevents surface water from flooding underground seams and contains spontaneous coal seam fires.",
                "Aatmanirbhar Bharat (Strategic): Fully indigenous, low-cost hardware and open-source software stack tailored for Indian mines."
            ])

    # ==========================================
    # SLIDE 6: RESEARCH, REFERENCES & FUTURE ROADMAP
    # ==========================================
    slide6 = prs.slides[5]
    for shape in slide6.shapes:
        if shape.name == "Title 1":
            shape.text_frame.text = "RESEARCH, REFERENCES & FINALE ROADMAP"
            shape.text_frame.paragraphs[0].font.size = Pt(24)
            shape.text_frame.paragraphs[0].font.bold = True
        elif shape.name == "TextBox 8":
            shape.text_frame.clear()
            tf = shape.text_frame
            tf.word_wrap = True

            def add_sec_6(title, bullets):
                p_head = tf.add_paragraph()
                p_head.space_before = Pt(4)
                p_head.space_after = Pt(2)
                run = p_head.add_run()
                run.text = title
                run.font.bold = True
                run.font.size = Pt(13)
                run.font.color.rgb = ACCENT_BLUE
                for b in bullets:
                    p = tf.add_paragraph()
                    p.space_after = Pt(2)
                    run_b = p.add_run()
                    run_b.text = "• " + b
                    run_b.font.size = Pt(11.5)
                    run_b.font.color.rgb = TEXT_DARK

            add_sec_6("1. 90-Day Grand Finale Execution Roadmap", [
                "LoRa Sub-GHz Mesh: Transition from ESP-NOW to SX1278 / SX1262 LoRa modules (Slotted TDMA) for 100+ nodes over a 5 km mining panel.",
                "Solar Deep-Sleep Power: 21700 4000mAh cell + 1.5W solar lid + 10µA deep-sleep -> 3–5+ years maintenance-free operation.",
                "Crack Extensometers: Spring-loaded draw-wire sensors across surface tension cracks for millimeter-level displacement tracking.",
                "Predictive AI Modeling: Knothe Time-Factor subsidence bowl model + small LSTM temporal forecast predicting ground sag 24–48 hours ahead.",
                "Industrial Edge Gateway: Raspberry Pi Gateway with 4G LTE module for offline pit-office GIS maps and emergency SMS broadcasting."
            ])
            add_sec_6("2. Core References & Technical Standards", [
                "Coal India Limited (CIL) SIH26025 Problem Statement Guidelines.",
                "DGMS Technical Circulars on Ground Control & Subsidence Management in Underground Coal Mines.",
                "CMPDI Standard Guidelines on Underground Strata & Surface Subsidence Prediction Curves.",
                "Liu et al., 'Isolation Forest for Unsupervised Anomaly Detection', IEEE ICDM.",
                "Madgwick / Mahony Sensor Fusion Algorithms for Low-Cost MEMS Attitude & Heading Reference Systems."
            ])

    # If there is a 7th instruction slide in the template, remove it so the final PPT is strictly 6 slides
    if len(prs.slides) > 6:
        rId = prs.slides._sldIdLst[6].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[6]

    prs.save(output_path)
    prs.save(clean_copy_path)
    prs.save(downloads_path)
    print("SUCCESS: Final 6-Slide presentation generated and saved to all locations!")

if __name__ == "__main__":
    build_presentation()
